from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title and Services Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_1.shapes.title
services = slide_1.placeholders[1]
title.text = "AI/ML Engineer\nMachine Learning Solutions"
services.text = "- Custom ML Models\n- Data Analysis\n- Deep Learning Algorithms\n- Model Training\n- Object Detection"

# Slide 2: Visual Representation Slide
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
visuals = slide_2.placeholders[1]
title.text = "Visual Representation"
visuals.text = "Include graphics depicting neural networks, data patterns, and machine learning algorithms."





# Save the presentation
presentation.save("AI_ML_Engineer_Gig.pptx")
